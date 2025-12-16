import pypdf
from pptx import Presentation

def extractPdfText(filePath):
    try:
        reader = pypdf.PdfReader(filePath)
        pages = []
        for page in reader.pages:
            content = page.extract_text()
            if content:
                pages.append(content)
        return "\n".join(pages)
    except Exception as e:
        print(f"PDF read error ({filePath}): {e}")
        return None

def extractPptxText(filePath):
    try:
        pres = Presentation(filePath)
        textParts = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    textParts.append(shape.text)
        return "\n".join(textParts)
    except Exception as e:
        print(f"PPTX read error ({filePath}): {e}")
        return None

def extractTxtText(filePath):
    try:
        with open(filePath, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        try:
            with open(filePath, 'r', encoding='latin-1') as f:
                return f.read()
        except Exception as e:
            print(f"TXT read error ({filePath}): {e}")
            return None
    except Exception as e:
        print(f"TXT read error ({filePath}): {e}")
        return None